"""
PDF to PowerPoint Converter Service.

Renders each PDF page to a full-bleed image and places it on its own slide.
This trades editability for reliable, high-fidelity visual output (LibreOffice's
PDF-to-PPTX conversion is generally poor and produces broken text boxes).
"""
import io
import os
from pathlib import Path
from typing import Optional, Tuple

import pymupdf
from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400


class PDFToPPTXConverterService:
    """Service for converting a PDF into an image-based PowerPoint deck"""

    def convert_pdf_to_pptx(self, pdf_path: str, output_path: str, dpi: int = 150) -> Tuple[bool, Optional[str]]:
        """
        Convert every page of a PDF into one slide, rendered as a full-slide image.

        Returns:
            Tuple of (success: bool, error_message: Optional[str])
        """
        try:
            if not os.path.exists(pdf_path):
                return False, f"PDF file not found: {pdf_path}"
            if os.path.getsize(pdf_path) == 0:
                return False, "PDF file is empty"

            output_dir = Path(output_path).parent
            output_dir.mkdir(parents=True, exist_ok=True)

            doc = pymupdf.open(pdf_path)
            try:
                if doc.page_count == 0:
                    return False, "PDF has no pages"

                prs = Presentation()
                blank_layout = prs.slide_layouts[6]

                for index in range(doc.page_count):
                    page = doc[index]
                    pix = page.get_pixmap(dpi=dpi)
                    png_bytes = pix.tobytes("png")

                    if index == 0:
                        page_width_in = page.rect.width / 72
                        page_height_in = page.rect.height / 72
                        prs.slide_width = Emu(int(page_width_in * EMU_PER_INCH))
                        prs.slide_height = Emu(int(page_height_in * EMU_PER_INCH))

                    slide = prs.slides.add_slide(blank_layout)
                    slide.shapes.add_picture(
                        io.BytesIO(png_bytes),
                        0,
                        0,
                        width=prs.slide_width,
                        height=prs.slide_height,
                    )
            finally:
                doc.close()

            prs.save(output_path)

            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                return False, "Failed to generate PowerPoint file"

            return True, None
        except Exception as e:
            return False, f"Unexpected error during PDF to PowerPoint conversion: {str(e)}"
